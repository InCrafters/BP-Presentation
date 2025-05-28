#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Criar uma nova apresentação
prs = Presentation()

# Definir cores da apresentação
COR_PRIMARIA = RGBColor(0, 82, 156)  # Azul escuro
COR_SECUNDARIA = RGBColor(0, 171, 169)  # Azul turquesa
COR_DESTAQUE = RGBColor(227, 30, 36)  # Vermelho
COR_TEXTO = RGBColor(51, 51, 51)  # Cinza escuro
COR_FUNDO = RGBColor(240, 240, 240)  # Cinza claro

# Slide 1: Capa
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Solução de Centralização de Cadastros de Usuários"
subtitle.text = "Hospital Beneficência Portuguesa"

# Ajustar formatação do título
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA
title.text_frame.paragraphs[0].font.bold = True

# Ajustar formatação do subtítulo
subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].font.color.rgb = COR_SECUNDARIA

# Slide 2: Sumário Executivo
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Sumário Executivo"
content_text = """
• Desafio: Gerenciamento de múltiplos cadastros de usuários em diferentes sistemas
• Problema: Redundância, inconsistências e dificuldades na gestão de fotos
• Solução: Centralização de cadastros para simplificar a gestão e garantir consistência
• Recomendação: Utilizar sistema base existente como fonte primária de cadastros
• Benefícios: Menor curva de aprendizado, aproveitamento de investimentos, implementação mais rápida
"""

content.text = content_text
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Slide 3: Contextualização do Problema
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Contextualização do Problema"
content_text = """
Desafios Atuais:
• Multiplicidade de Cadastros em diferentes sistemas (W-Access, MV, RH e outros)
• Inconsistência de Dados entre sistemas para o mesmo usuário
• Gestão Complexa de Fotos dos usuários
• Retrabalho Operacional no cadastro de usuários
• Experiência Fragmentada para operadores e usuários finais

Impactos:
• Aumento do tempo em tarefas administrativas
• Maior probabilidade de erros humanos
• Dificuldades no controle de acesso físico e lógico
• Comprometimento da segurança da informação
"""

content.text = content_text
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Slide 4: Cenário 1 (Recomendado)
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout com título e conteúdo
title = slide.shapes.title
title.text = "Cenário 1: Sistema Base como Fonte Primária (Recomendado)"
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Adicionar imagem do diagrama
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4.5)
slide.shapes.add_picture("/home/ubuntu/apresentacao_bp/imagens/diagrama_cenario1.png", left, top, width, height)

# Slide 5: Vantagens do Cenário 1
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Vantagens do Cenário 1 (Recomendado)"
content_text = """
• Curva de aprendizado reduzida: Operadores já familiarizados com o sistema base
• Aproveitamento de investimento existente: Utiliza infraestrutura já adquirida
• Menor resistência à mudança: Usuários continuam usando sistema conhecido
• Implementação mais rápida: Apenas desenvolvimento de conectores de integração
• Menor risco operacional: Sistema base já em produção com estabilidade comprovada
• Continuidade operacional: Sem interrupção nos processos durante implementação
• Custo-benefício otimizado: Melhor relação entre investimento e resultados
"""

content.text = content_text
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Slide 6: Cenário 2
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Cenário 2: Nova Aplicação Centralizada para Cadastros"
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Adicionar imagem do diagrama
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4.5)
slide.shapes.add_picture("/home/ubuntu/apresentacao_bp/imagens/diagrama_cenario2.png", left, top, width, height)

# Slide 7: Cenário 3
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Cenário 3: Aplicação como Serviço via API"
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Adicionar imagem do diagrama
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4.5)
slide.shapes.add_picture("/home/ubuntu/apresentacao_bp/imagens/diagrama_cenario3.png", left, top, width, height)

# Slide 8: Análise Comparativa
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Análise Comparativa dos Cenários"
content_text = """
Critérios de Avaliação:

• Tempo de implementação: Cenário 1 (Menor) < Cenário 3 (Médio) < Cenário 2 (Maior)
• Custo: Cenário 1 (Menor) < Cenário 3 (Médio) < Cenário 2 (Maior)
• Curva de aprendizado: Cenário 1 (Mínima) < Cenário 3 (Média) < Cenário 2 (Alta)
• Flexibilidade técnica: Cenário 1 (Média) < Cenário 2 (Alta) = Cenário 3 (Alta)
• Risco operacional: Cenário 1 (Baixo) < Cenário 3 (Médio) < Cenário 2 (Alto)
• Manutenção a longo prazo: Cenário 1 (Média) = Cenário 2 (Média) < Cenário 3 (Alta)
• Escalabilidade: Cenário 1 (Média) < Cenário 2 (Alta) = Cenário 3 (Alta)

Conclusão: O Cenário 1 apresenta o melhor equilíbrio entre benefícios, custos e facilidade de implementação.
"""

content.text = content_text
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Slide 9: Plano de Implementação
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Plano de Implementação"
content_text = """
Fase 1: Preparação e Análise (4-6 semanas)
• Mapeamento detalhado dos sistemas e interfaces
• Definição do sistema base a ser utilizado
• Análise de requisitos técnicos e funcionais

Fase 2: Desenvolvimento (8-12 semanas)
• Implementação do middleware de integração
• Desenvolvimento dos conectores específicos
• Criação do repositório centralizado de fotos

Fase 3: Implantação (4-6 semanas)
• Testes de aceitação com usuários-chave
• Migração e sincronização inicial de dados
• Implantação gradual por departamentos

Fase 4: Monitoramento e Otimização (Contínuo)
• Acompanhamento de métricas de desempenho
• Ajustes e otimizações conforme necessário
"""

content.text = content_text
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Slide 10: Benefícios Esperados
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Benefícios Esperados"
content_text = """
• Consistência de Dados: Garantia de que todos os sistemas possuem as mesmas informações
• Simplificação da Gestão de Fotos: Processo unificado para atualização de imagens
• Redução de Retrabalho: Eliminação da necessidade de múltiplos cadastros
• Aumento de Produtividade: Menos tempo dedicado a tarefas administrativas
• Melhoria na Segurança: Controle centralizado de acessos e permissões

Retorno Esperado:
• Redução de 70% no tempo dedicado à gestão de cadastros
• Eliminação de inconsistências entre sistemas
• Aumento na segurança e controle de acesso
• Melhoria na experiência dos usuários
• Redução de custos operacionais a longo prazo
"""

content.text = content_text
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Slide 11: Próximos Passos
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Próximos Passos"
content_text = """
1. Workshop de Detalhamento
   • Sessão com equipes técnicas e operacionais para refinamento da solução

2. Prova de Conceito
   • Implementação piloto com escopo reduzido para validação técnica

3. Planejamento Detalhado
   • Cronograma, recursos e marcos para implementação completa

4. Início do Projeto
   • Mobilização da equipe e início das atividades
"""

content.text = content_text
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Slide 12: Conclusão
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusão"
content_text = """
A centralização de cadastros de usuários representa uma oportunidade significativa para o Hospital Beneficência Portuguesa:

• Otimizar processos administrativos
• Reduzir custos operacionais
• Melhorar a experiência de colaboradores e pacientes
• Aumentar a segurança da informação
• Garantir consistência de dados em todos os sistemas

A solução recomendada (Cenário 1) oferece o melhor equilíbrio entre benefícios, custos e facilidade de implementação, permitindo resultados rápidos com menor impacto operacional.

Estamos à disposição para esclarecer dúvidas e avançar com os próximos passos para a implementação desta solução transformadora.
"""

content.text = content_text
title.text_frame.paragraphs[0].font.color.rgb = COR_PRIMARIA

# Salvar a apresentação
prs.save('/home/ubuntu/apresentacao_bp/Apresentacao_Centralizacao_Cadastros_BP.pptx')
